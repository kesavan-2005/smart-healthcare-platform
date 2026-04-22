import sys
import subprocess

def install(package):
    subprocess.check_call([sys.executable, "-m", "pip", "install", package])

try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    install("python-pptx")
    import pptx

def extract_text_from_pptx(path):
    print(f"Reading {path}...")
    prs = pptx.Presentation(path)
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"\n--- SLIDE {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    with open('d:/smart-healthcare-platform/ppt_text.txt', 'w', encoding='utf-8') as f:
        f.write(text)
    print("Extracted text saved to ppt_text.txt")
    print(text[:2000])

extract_text_from_pptx('d:/smart-healthcare-platform/2nd  review.pptx')
